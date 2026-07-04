"""
tools/executor.py — Hardened Executor v14.0 (Async + NATO Vocal MFA).

Security layers:
  Layer 1 — Allowlist + flag blocking: only approved binaries; rejects
             -EncodedCommand (PS) and python -c (inline eval).
  Layer 2 — Path canonicalization: Path.resolve() on every token; blocks
             paths under C:/Windows, System32, or root /.
  Layer 3 — Async NATO OTP: _challenge() awaits an asyncio.Queue fed by a
             background STT thread via loop.call_soon_threadsafe. 30s timeout.
  Layer 4 — Execution in run_in_executor: subprocess.run never blocks the
             event loop.
  Layer 5 — shell=False always. No exceptions.
  Layer 6 — Regex validation on all network inputs (domains, IPs).
"""

import asyncio
import ipaddress
import json
import random
import re
import sys
import os
import shlex
import subprocess
import platform
import shutil
import http.server
import socketserver
import threading
import webbrowser
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
import psutil
import requests
from loguru import logger
import whois
import dns.resolver

from core.ironman_mode import SessionConsent, default_consent
from core.risk_classes import (
    classify_tool,
    requires_hitl,
    requires_trusted_lab,
    rollback_hint,
    binary_risk_class,
    verify_consistent_with_legacy_sets,
)

# ── NATO Vocal MFA ───────────────────────────────────────────────────────────
_NATO_ALPHABET: tuple[str, ...] = (
    "Alpha", "Bravo", "Charlie", "Delta", "Echo", "Foxtrot",
    "Golf", "Hotel", "India", "Juliet", "Kilo", "Lima",
    "Mike", "November", "Oscar", "Papa", "Quebec", "Romeo",
    "Sierra", "Tango", "Uniform", "Victor", "Whiskey", "Xray",
    "Yankee", "Zulu",
)
_VOCAL_CONFIRM_WORDS: frozenset[str] = frozenset({
    "hazlo", "sí", "si", "yes", "confirmar", "autorizar", "ejecutar", "adelante", "execute",
})
_CONFIDENCE_THRESHOLD = 0.75

# ── Layer 1: Allowlist de ejecutables permitidos ──────────────────────────────
COMMAND_ALLOWLIST: frozenset[str] = frozenset({
    # Networking (diagnóstico / reconocimiento)
    "ping", "nmap", "whois", "traceroute", "tracert",
    "netstat", "ipconfig", "ifconfig", "arp",
    "curl", "wget",
    # Información del sistema (lectura)
    "ps", "top", "htop", "tasklist",
    "df", "du", "free",
    "uname", "hostname", "whoami", "id",
    # Dev tools
    "python", "python3", "pip", "pip3",
    "git", "node", "npm",
    "gcc", "make",
    # Navegación de archivos (lectura)
    "ls", "dir", "cat", "type", "more",
    "grep", "find", "findstr",
    "head", "tail", "wc",
    # Miscelánea segura
    "echo", "ssh", "scp", "openssl",
})

# Lab-only binaries: offensive / heavier tooling permitted through
# RedTeamShellExecutor ONLY when trusted-lab mode is explicitly enabled, and
# always at FULL_NATO (they are never in COMMAND_ALLOWLIST, so _classify marks
# them unlisted_escalation → FULL_NATO). This is an explicit allowlist —
# trusted-lab never means "run any binary".
_LAB_COMMAND_ALLOWLIST: frozenset[str] = frozenset({
    "masscan", "nikto", "hydra", "sqlmap", "gobuster", "ffuf", "dirb",
    "msfconsole", "msfvenom", "sliver", "tcpdump", "tshark",
    "hashcat", "john", "responder", "crackmapexec",
})

# ── RedTeamShellExecutor: hard-block patterns (no OTP override) ───────────────
_CRITICAL_BLOCK: list[str] = [
    "rm -rf /",
    "del /f /s /q c:\\windows",
    "format c:",
    "dd if=/dev/zero",
    "> /dev/sda",
]

# Patterns that elevate a command to requires_challenge=True
_SUSPICIOUS_PATTERNS: list[str] = [
    "-encodedcommand", "-enc",
    "invoke-expression", "iex",
    "virtualalloc", "writeprocessmemory",
    "net user", "net localgroup",
]

# ── Layer 1: Flags bloqueados (evasión de policy) ────────────────────────────
_BLOCKED_FLAGS: frozenset[str] = frozenset({
    "-encodedcommand", "-enc",  # PowerShell encoded command execution
    "--encoded",
})
_PYTHON_EXECUTABLES: frozenset[str] = frozenset({"python", "python3"})

# ── Layer 2: Directorios del sistema bloqueados ───────────────────────────────
def _build_system_dirs() -> frozenset[Path]:
    dirs: set[Path] = {Path("/").resolve()}
    if os.name == "nt":
        sysroot = Path(os.environ.get("SystemRoot", r"C:\Windows"))
        dirs.add(sysroot.resolve())
        dirs.add((sysroot / "System32").resolve())
        dirs.add(Path(r"C:\Windows").resolve())
        dirs.add(Path(r"C:\Windows\System32").resolve())
    else:
        dirs.update({
            Path("/etc").resolve(),
            Path("/sys").resolve(),
            Path("/proc").resolve(),
            Path("/boot").resolve(),
        })
    return frozenset(dirs)

_SYSTEM_DIRS: frozenset[Path] = _build_system_dirs()

# Metacaracteres de shell: si aparecen en el input, el comando se rechaza
_FORBIDDEN_CHARS_RE = re.compile(r'[;&|`$<>()\{\}!\\\n\r]')

# Regex para hosts válidos: IP, CIDR, hostname (sin metacaracteres)
_SAFE_HOST_RE = re.compile(
    r'^(?:[a-zA-Z0-9](?:[a-zA-Z0-9\-]{0,61}[a-zA-Z0-9])?\.)*'
    r'[a-zA-Z0-9](?:[a-zA-Z0-9\-]{0,61}[a-zA-Z0-9])?'
    r'(?:/\d{1,2})?$'
)

# Regex para dominios: solo alfanuméricos, puntos y guiones
_SAFE_DOMAIN_RE = re.compile(r'^[a-zA-Z0-9._-]{1,253}$')

_CSP_META = (
    '<meta http-equiv="Content-Security-Policy" '
    "content=\"default-src 'self' 'unsafe-inline'; "
    "img-src 'self' data:; connect-src 'none';\">"
)

# Tools cuya seguridad es read-only/benign — no requieren desafío HITL.
# run_shell_command fue removido: ahora pasa por el desafío NATO vocal.
_HITL_EXEMPT_TOOLS: frozenset[str] = frozenset({
    "get_datetime",
    "get_weather",
    "web_search",
    "fetch_webpage",
    "system_info",
    "list_processes",
    "list_directory",
    "read_file",
    "leer_archivo_universal",
    "escanear_pantalla",
    "analizar_codigo_sast",
    "get_clipboard",
    "check_connectivity",
    "whois_lookup",
    "consultar_base_conocimiento",
    "estudiar_tema",
    "get_system_status",
    "ingest_docs",
    "query_knowledge",
    # V59.0 APEX — new read-only / safe tools
    "decode_payload",
    "hash_file",
    "port_lookup",
    "regex_test",
    "list_notes",
    "git_query",
    "save_note",
})

# Tools that ALWAYS require explicit HITL/NATO approval — arbitrary code / shell
# / outbound-network execution surfaces. They can NEVER be HITL-exempt and can
# never be auto-approved by any trust path (F4). Even if one is mistakenly added
# to _HITL_EXEMPT_TOOLS, the aexecute gate still forces a challenge for it.
_ALWAYS_HITL_TOOLS: frozenset[str] = frozenset({
    "code_execute",
    "run_shell_command",
    "http_request",
})

# Security invariant: no always-HITL tool may ever appear in the exempt list.
assert _ALWAYS_HITL_TOOLS.isdisjoint(_HITL_EXEMPT_TOOLS), (
    "SECURITY: an always-HITL tool is present in _HITL_EXEMPT_TOOLS"
)

# V62.0 Phase 7 — the risk-class taxonomy (core/risk_classes.py) is now the
# live gating decision in aexecute()/aexecute_mcp() below. This assertion
# guarantees it can never silently diverge from the two legacy sets above
# (which 5 other test files assert on directly) — a mismatch here is a
# security bug, caught at import time, not at runtime.
verify_consistent_with_legacy_sets(_HITL_EXEMPT_TOOLS, _ALWAYS_HITL_TOOLS)

# ── MCP gateway ────────────────────────────────────────────────────────────
# Tools proxied from an external MCP server (e.g. the packet_tracer_bridge
# stdio bridge) are foreign, unaudited code — they must never bypass this
# module's security gate. Allowlist-not-denylist: an MCP tool name that isn't
# listed here is refused outright, and every allowlisted MCP tool is ALWAYS
# HITL-challenged (no exempt tier, unlike local tools — see aexecute_mcp).
MCP_TOOL_ALLOWLIST: frozenset[str] = frozenset({
    "generar_laboratorio_red",
    "abrir_packet_tracer",
})

# Argument keys whose value names a file: must resolve to a bare basename
# with no path separators, drive letter, or traversal segments.
_MCP_FILENAME_ARG_KEYS: frozenset[str] = frozenset({"nombre_archivo"})


def _validate_mcp_filename(value: str) -> str | None:
    """Reject anything that isn't a safe bare filename.

    Returns a human-readable error string, or None if *value* is safe to
    join onto a fixed base directory. Catches path traversal ('../x'),
    absolute paths, and drive-letter/UNC prefixes by comparing against
    Path(value).name — any of those makes the basename differ from the
    original string.
    """
    if not value or not isinstance(value, str):
        return "Nombre de archivo vacío o inválido."
    if "\x00" in value:
        return "Nombre de archivo contiene un byte nulo."
    if len(value) > 255:
        return "Nombre de archivo demasiado largo."
    if value in (".", "..") or Path(value).name != value:
        return f"Nombre de archivo inválido (path traversal o separadores no permitidos): '{value}'"
    return None

# Patrones SAST precompilados para _tool_analizar_codigo_sast
_SAST_PATTERNS: list[tuple] = [
    (re.compile(r'eval\('),                                            "eval() detectado"),
    (re.compile(r'exec\('),                                            "exec() detectado"),
    (re.compile(r'os\.system\('),                                      "os.system() detectado"),
    (re.compile(r'(?i)(password|api_key|secret|token)\s*=\s*["\']'),  "Credencial hardcodeada"),
    (re.compile(r'(?i)SELECT\s+.+\s+WHERE\s+.+=\s*%s'),               "Concatenación SQL raw"),
]

# Map app_name → executable: el input del usuario solo se usa como clave de lookup,
# nunca se interpola en el comando. shell=False siempre.
_APP_SOFTWARE_MAP: dict[str, str] = {
    "word": "winword",
    "excel": "excel",
    "powerpoint": "powerpnt",
    "access": "msaccess",
    "outlook": "outlook",
    "autocad": "acad",
    "packet tracer": "PacketTracer",
    "blender": "blender",
    "chrome": "chrome",
    "firefox": "firefox",
    "vscode": "code",
    "notepad": "notepad",
    "calculator": "calc",
    "paint": "mspaint",
    "vlc": "vlc",
    "obs": "obs64",
    "spotify": "Spotify",
    "discord": "Discord",
    "teams": "Teams",
    "zoom": "Zoom",
    "gimp": "gimp-2.10",
    "inkscape": "inkscape",
    "wireshark": "Wireshark",
    "burpsuite": "burpsuite",
    "virtualbox": "VirtualBox",
    "vmware": "vmware",
}


# ── Guardrail: forbidden destructive patterns on system directories ───────────
_GUARDRAIL_SYSTEM_WRITE_RE = re.compile(
    r'(?i)(?:'
    r'(?:del|rm|erase|rd|rmdir)\s+.*?(?:C:[/\\]Windows|System32)'
    r'|(?:copy|move|xcopy|robocopy|cp|mv)\s+.*?\s+(?:C:[/\\]Windows|System32)'
    r'|reg\s+(?:add|delete|import)\b'
    r')'
)

_GUARDRAIL_ROOT_DELETE_RE = re.compile(
    r'(?i)(?:'
    r'rm\s+-[frRFR]+\s+[/\\]\s*(?:\s|$)'
    r'|del\s+(?:/[sqSQ]\s+)*C:[/\\]\s*(?:\s|$)'
    r'|rd\s+(?:/[sqSQ]\s+)*C:[/\\]\s*(?:\s|$)'
    r'|rmdir\s+(?:/[sqSQ]\s+)*C:[/\\]\s*(?:\s|$)'
    r'|format\s+[cCdD]:'
    r')'
)

# ── PII detection in tool outputs ─────────────────────────────────────────────
_PII_OUTPUT_RE = re.compile(
    r'(?i)(?:'
    r'(?:password|passwd|contraseña)\s*[:=]\s*\S{4,}'
    r'|(?:api[_-]?key|secret[_-]?key|access[_-]?key)\s*[:=]\s*[A-Za-z0-9_\-]{16,}'
    r'|(?:bearer\s+token|auth[_-]?token)\s*[:=]\s*[A-Za-z0-9_\-\.]{20,}'
    r')'
)


def _contains_shell_metacharacters(value: str) -> bool:
    """True if *value* contains any shell metacharacter / expansion token."""
    return bool(_FORBIDDEN_CHARS_RE.search(value or ""))


# ── Centralized filesystem sandbox ────────────────────────────────────────────
# Allowed write/read roots, resolved once at call time. Mirrors the inline
# containment check that _tool_read_file / _tool_write_file already perform, so
# every path-taking handler can share one hardened definition instead of
# re-implementing (and drifting from) it.
def _sandbox_allowed_dirs() -> tuple[Path, ...]:
    return (
        Path.home() / "Downloads",
        Path.home() / "Documents",
        Path.cwd(),
    )


def _resolve_within_allowed(path: str) -> "Path | None":
    """Resolve *path* and return it iff it is contained within an allowed dir.

    Returns ``None`` — i.e. fail-closed — on every escape attempt: relative
    traversal (``../``), absolute paths outside the roots, drive-letter escapes,
    or symlinks whose target lands outside (``.resolve()`` follows symlinks and
    normalizes ``..`` *before* the containment test, so a symlink inside an
    allowed dir pointing outside is still rejected). Any resolution failure
    (malformed path, OS error) is likewise treated as not-allowed.
    """
    if not isinstance(path, str) or not path.strip():
        return None
    try:
        p = Path(path).expanduser().resolve()
    except (OSError, ValueError, RuntimeError):
        return None
    for allowed in _sandbox_allowed_dirs():
        try:
            if p.is_relative_to(allowed.resolve()):
                return p
        except (OSError, ValueError, RuntimeError):
            continue
    return None


def _strip_override(tool_name: str, tool_input: dict) -> dict:
    """Return a copy of *tool_input* with any model-supplied FORCE_OVERRIDE removed.

    Defense-in-depth: the LLM must never be able to pass a guardrail-bypass flag
    into a handler. If present, it is logged as a probable injection attempt and
    dropped before validation or execution.
    """
    if not isinstance(tool_input, dict) or "FORCE_OVERRIDE" not in tool_input:
        return tool_input
    logger.warning(
        f"SECURITY: stripped model-supplied FORCE_OVERRIDE from tool='{tool_name}' "
        "input (guardrail bypass via tool argument is disabled)."
    )
    return {k: v for k, v in tool_input.items() if k != "FORCE_OVERRIDE"}


def _trusted_lab_enabled() -> bool:
    """True only when the operator has explicitly enabled trusted-lab mode.

    Read dynamically (config first, raw env fallback) so the value is never
    cached from a model-generated argument and can be toggled per-process for
    tests. This is the ONLY legitimate source of a security override — tool
    input is never consulted.
    """
    try:
        from core.config import settings
        if settings.trusted_lab_mode:
            return True
    except Exception:
        pass
    return os.environ.get("JARVIS_TRUSTED_LAB", "").strip().lower() in {
        "1", "true", "yes", "on",
    }


def _http_target_blocked(url: str) -> str | None:
    """SSRF guard for outbound HTTP tools.

    Resolves the URL host and rejects loopback, RFC1918 private, link-local
    (incl. 169.254.169.254 cloud metadata), unique-local, multicast, and other
    reserved ranges — unless trusted-lab mode is explicitly enabled. Returns a
    human-readable block reason, or None if the target is permitted.
    """
    import socket
    import urllib.parse

    if _trusted_lab_enabled():
        return None

    parsed = urllib.parse.urlparse(url)
    if parsed.scheme not in ("http", "https"):
        return f"Esquema no permitido: {parsed.scheme or '(vacío)'} (usa http/https)."
    host = parsed.hostname
    if not host:
        return "URL inválida: host vacío."

    # Resolve every address the host maps to; block if ANY is internal — this
    # defeats DNS-rebinding and hostnames that alias a private/metadata IP.
    candidates: set[str] = set()
    try:
        ipaddress.ip_address(host)
        candidates.add(host)
    except ValueError:
        try:
            for fam, _, _, _, sockaddr in socket.getaddrinfo(host, None):
                candidates.add(sockaddr[0])
        except Exception:
            return f"No se pudo resolver el host '{host}'."

    for raw_ip in candidates:
        try:
            ip = ipaddress.ip_address(raw_ip.split("%")[0])  # strip zone id
        except ValueError:
            return f"Dirección IP no válida resuelta para '{host}'."
        if (
            ip.is_private or ip.is_loopback or ip.is_link_local
            or ip.is_multicast or ip.is_reserved or ip.is_unspecified
        ):
            return (
                f"Destino interno bloqueado (SSRF): {host} → {ip}. "
                "Habilita JARVIS_TRUSTED_LAB=true para permitir rangos internos en lab aislado."
            )
    return None


def _validate_network_target(target: str) -> str:
    """
    Validate a scan/connectivity target. Returns the normalized target string.

    Accepts:
      - IPv4 addresses                (192.168.1.1)
      - IPv4 CIDR ranges, prefix 0-32 (10.0.0.0/24)
      - safe hostnames / domains      (scanme.nmap.org)

    Rejects empty values, whitespace, path traversal, URL schemes and any
    shell-injection payload. Raises ValueError on an invalid target.
    """
    t = (target or "").strip()
    if not t:
        raise ValueError("Target inválido: valor vacío.")
    if " " in t or _contains_shell_metacharacters(t):
        raise ValueError("Target inválido. Usa una IP, rango CIDR o hostname válido.")
    # IP / CIDR first — ipaddress enforces sane prefix bounds (0-32).
    try:
        if "/" in t:
            ipaddress.ip_network(t, strict=False)
        else:
            ipaddress.ip_address(t)
        return t
    except ValueError:
        pass
    # Fall back to a conservative hostname check (no scheme, no traversal).
    if _SAFE_HOST_RE.match(t):
        return t
    raise ValueError("Target inválido. Usa una IP, rango CIDR o hostname válido.")


def _validate_command(
    command: str, extra_allowlist: frozenset[str] = frozenset()
) -> tuple[bool, str, list[str]]:
    """
    Valida y parsea un comando contra la allowlist (Layers 1 & 2).

    ``extra_allowlist`` adds binaries to the base COMMAND_ALLOWLIST for this call
    only (used for explicit lab binaries under trusted-lab mode). It defaults to
    empty, so normal callers keep the exact base-allowlist behavior.

    Returns:
        (is_valid, error_message, argv)
        Si is_valid es False, argv es [].
    """
    if not command.strip():
        return False, "Comando vacío.", []

    # Layer 1a: Bloquear metacaracteres de shell antes de parsear
    if _FORBIDDEN_CHARS_RE.search(command):
        return (
            False,
            "Comando rechazado: contiene metacaracteres de shell prohibidos "
            f"({_FORBIDDEN_CHARS_RE.pattern}).",
            [],
        )

    # Layer 1b: Parseo seguro con shlex
    try:
        argv = shlex.split(command)
    except ValueError as e:
        return False, f"Comando malformado: {e}", []

    if not argv:
        return False, "Comando vacío tras parseo.", []

    # Layer 1c: Normalizar el nombre del ejecutable (quitar ruta y .exe en Windows)
    executable = Path(argv[0]).name.lower().removesuffix(".exe")

    # Layer 1d: Verificar contra la allowlist (+ optional lab allowlist)
    allowed = COMMAND_ALLOWLIST | extra_allowlist
    if executable not in allowed:
        return (
            False,
            f"Ejecutable '{executable}' no está en la allowlist.\n"
            f"Permitidos: {', '.join(sorted(allowed))}",
            [],
        )

    # Layer 1e: Bloquear flags de evasión (-EncodedCommand, python -c, etc.)
    for arg in argv[1:]:
        if arg.lower() in _BLOCKED_FLAGS:
            return (
                False,
                f"Flag '{arg}' bloqueado — evasión de política de ejecución.",
                [],
            )
        if executable in _PYTHON_EXECUTABLES and arg == "-c":
            return (
                False,
                "python -c bloqueado — ejecución de código inline no permitida.",
                [],
            )

    # Layer 2: Canonicalización de rutas — bloquear acceso a directorios del sistema
    for token in argv[1:]:
        is_path_like = "/" in token or "\\" in token or (
            len(token) >= 3 and token[1:3] in (":/", ":\\")
        )
        if not is_path_like:
            continue
        try:
            resolved = Path(token).resolve()
            for sys_dir in _SYSTEM_DIRS:
                if resolved == sys_dir or sys_dir in resolved.parents:
                    return (
                        False,
                        f"Ruta bloqueada: '{resolved}' apunta a un directorio del sistema.",
                        [],
                    )
        except Exception:
            pass

    return True, "", argv


async def _aura_broadcast(event: dict) -> None:
    """Fire-and-forget telemetry to the AURA WebSocket pipeline.

    Imported lazily so executor.py has no hard dependency on the AURA module —
    Jarvis works in text mode even when FastAPI/uvicorn are not installed.
    """
    try:
        from aura.server import broadcast
        await broadcast(event)
    except Exception:
        pass


class ToolExecutor:
    def __init__(
        self,
        stt_queue: "asyncio.Queue | None" = None,
        stt_listener=None,
        consent: "SessionConsent | None" = None,
    ) -> None:
        self._active_web_server: socketserver.TCPServer | None = None
        self._active_web_thread: threading.Thread | None = None
        self._stt_queue = stt_queue        # asyncio.Queue[(str, float)] | None
        self._stt_listener = stt_listener  # HighPrioritySTTListener | None
        # V62.0 Phase 6 — consent-gated sensitive surfaces. Defaults fully OFF
        # (fail-closed) when the caller doesn't wire a shared session consent —
        # see core.ironman_mode.SessionConsent / core.consent_commands.
        self.consent: SessionConsent = consent if consent is not None else default_consent()
        from core.governance import TacticAuditLogger
        self._audit = TacticAuditLogger()

    def _consent_error(self, surface: str, grant_phrase: str) -> dict:
        return {
            "error": (
                f"Consent required for {surface} — not granted this session. "
                f"Say '{grant_phrase}' to allow it."
            )
        }

    # ── Layer 3: NATO Vocal MFA ───────────────────────────────────────────────

    def _evaluate_nato_response(
        self, text: str, confidence: float, challenge_word: str
    ) -> bool:
        """Return True if the spoken response matches the NATO challenge word."""
        if confidence < _CONFIDENCE_THRESHOLD:
            logger.warning(f"VAP NATO: confianza baja ({confidence:.2%}) — denegado.")
            return False
        words = set(
            text.lower()
            .replace(",", "").replace(".", "").replace("¡", "").replace("!", "")
            .split()
        )
        if challenge_word.lower() in words:
            logger.info(f"VAP NATO: '{challenge_word}' confirmado (conf={confidence:.2%}).")
            return True
        if words & _VOCAL_CONFIRM_WORDS:
            logger.info(f"VAP: Confirmación general (conf={confidence:.2%}).")
            return True
        logger.warning(f"VAP NATO: Respuesta inválida '{text}' para '{challenge_word}'.")
        return False

    async def _challenge(self, tool_name: str, preview: str) -> tuple[bool, str]:
        """
        Layer 3 — Async NATO OTP gate.

        Displays a random NATO word challenge, starts a background audio-capture
        thread that pushes (text, confidence) into self._stt_queue via
        loop.call_soon_threadsafe, then awaits the result with a 30-second
        timeout. Falls back to keyboard via run_in_executor on timeout or when
        STT is unavailable.

        Returns:
            (granted: bool, auth_audit: str)
        """
        challenge_word = random.choice(_NATO_ALPHABET)
        bar = "=" * 62
        print(f"\n{bar}")
        print("  [!] AUTORIZACIÓN NATO REQUERIDA")
        print(f"      Tool      : {tool_name.upper()}")
        print(f"      Parámetros: {preview}")
        print(f"  >> DESAFÍO: Di la palabra NATO [{challenge_word.upper()}] para autorizar <<")
        print("  (o presiona 'y' en el teclado | timeout: 30s)")
        print(f"{bar}")
        sys.stdout.flush()

        loop = asyncio.get_running_loop()

        # Broadcast the challenge so the AURA UI can display it
        await _aura_broadcast({
            "type": "challenge",
            "tool": tool_name,
            "challenge_word": challenge_word,
            "timestamp": datetime.now(timezone.utc).isoformat(),
        })

        if self._stt_queue is not None and self._stt_listener is not None:
            def _capture_and_push() -> None:
                try:
                    audio = self._stt_listener.record()
                    text, confidence = self._stt_listener.transcribe_with_confidence(audio)
                    loop.call_soon_threadsafe(
                        self._stt_queue.put_nowait, (text, confidence)
                    )
                except Exception as exc:
                    logger.warning(f"VAP: Error de captura — {exc}")
                    loop.call_soon_threadsafe(
                        self._stt_queue.put_nowait, ("", 0.0)
                    )

            threading.Thread(
                target=_capture_and_push, daemon=True, name="vap-capture"
            ).start()

            try:
                text, confidence = await asyncio.wait_for(
                    self._stt_queue.get(), timeout=30.0
                )
                # Empty transcript or sub-threshold confidence is NOT a denial —
                # the operator may simply be in a noisy room or STT misfired.
                # Fall back to keyboard confirmation instead of hard-denying.
                if (not text or not text.strip()) or confidence < _CONFIDENCE_THRESHOLD:
                    logger.warning("VAP NATO: low confidence — falling back to keyboard")
                else:
                    granted = self._evaluate_nato_response(text, confidence, challenge_word)
                    auth_audit = (
                        f"vocal:nato:{challenge_word}:{confidence:.2f}:"
                        f"{'granted' if granted else 'denied'}"
                    )
                    if not granted:
                        logger.warning(f"VAP NATO: Denegado — texto='{text}'")
                    return granted, auth_audit
            except asyncio.TimeoutError:
                logger.info("VAP: Timeout vocal — fallback a teclado.")

        # Keyboard fallback — runs in thread pool so event loop stays free.
        # Reached on: STT unavailable, vocal timeout, or low-confidence transcript.
        # Security: still requires an explicit 'y'; never auto-approves.
        auth = await loop.run_in_executor(
            None, lambda: input("  ¿Autorizar ejecución? (y/N): ")
        )
        granted = auth.strip().lower() == "y"
        return granted, f"keyboard:{'granted' if granted else 'denied'}"

    # ── Async executor gate ───────────────────────────────────────────────────

    def _preflight_validate(
        self, tool_name: str, tool_input: dict, include_shell: bool = False
    ) -> dict | None:
        """
        Tool-specific input validation that runs *before* the generic
        destructive-pattern guardrail, so a precise '…inválido' message wins
        over the broad guardrail net for injection payloads embedded in a
        target/domain (e.g. "192.168.1.1; rm -rf /").

        ``include_shell`` adds allowlist/metachar validation for
        run_shell_command. The async path (aexecute) leaves it False so the
        handler's static-triage / neutralization pipeline still runs on a
        blocked command after the NATO gate.

        Returns an error dict if the input is rejected, else None.
        """
        if include_shell and tool_name == "run_shell_command":
            is_valid, error_msg, _ = _validate_command(str(tool_input.get("command", "")))
            if not is_valid:
                return {"error": f"Comando bloqueado por política de seguridad: {error_msg}"}
        if tool_name in ("network_scan", "check_connectivity"):
            field = "target" if tool_name == "network_scan" else "host"
            try:
                _validate_network_target(str(tool_input.get(field, "")))
            except ValueError as e:
                return {"error": str(e)}
            if tool_name == "network_scan":
                scan_type = str(tool_input.get("scan_type", "-sS -sV"))
                if _contains_shell_metacharacters(scan_type):
                    return {
                        "error": "scan_type inválido: contiene metacaracteres de shell prohibidos."
                    }
        elif tool_name == "whois_lookup":
            if not _SAFE_DOMAIN_RE.match(str(tool_input.get("domain", ""))):
                return {
                    "error": "Dominio inválido. Solo alfanuméricos, puntos y guiones (1-253 chars)."
                }
        return None

    def execute(self, tool_name: str, tool_input: dict, reasoning: str = "") -> dict:
        """
        Synchronous execution gate — the reusable, testable public API.

        Enforces every non-interactive security layer (tool-specific input
        validation, destructive-pattern guardrails, allowlist + shell=False
        inside each handler) and runs the handler in the calling thread.

        It deliberately does NOT run the interactive NATO vocal challenge:
        that Human-In-The-Loop gate lives in aexecute(), the async path the
        live orchestrator uses. Use aexecute() whenever you are inside the
        event loop and want the full authorization pipeline.
        """
        tool_input = _strip_override(tool_name, tool_input)
        handler = getattr(self, f"_tool_{tool_name}", None)
        if handler is None:
            self._audit.log_action(
                tool_name, reasoning, "unknown", "error", "Tool no implementada"
            )
            return {"error": f"Tool '{tool_name}' no implementada."}

        preflight = self._preflight_validate(tool_name, tool_input, include_shell=True)
        if preflight is not None:
            self._audit.log_action(
                tool_name, reasoning, "blocked:preflight", "blocked",
                preflight.get("error", "")[:200],
            )
            return preflight

        guardrail_block = self._validate_guardrails(tool_name, tool_input)
        if guardrail_block:
            self._audit.log_action(
                tool_name, reasoning, "blocked:guardrail", "blocked",
                guardrail_block.get("error", "")[:200],
            )
            return guardrail_block

        try:
            result = handler(**tool_input)
        except Exception as e:
            logger.error(f"Error en tool '{tool_name}': {e}")
            self._audit.log_action(tool_name, reasoning, "sync", "error", str(e)[:200])
            return {"error": str(e)}

        status = "error" if isinstance(result, dict) and "error" in result else "success"
        result = self._check_pii_output(result)
        self._audit.log_action(
            tool_name, reasoning, "sync", status,
            json.dumps(result, ensure_ascii=False, default=str)[:200],
        )
        return result

    async def aexecute(self, tool_name: str, tool_input: dict, reasoning: str = "") -> Any:
        """
        Fully async execution gate:
          1. Look up handler.
          2. Tool-specific pre-flight validation.
          3. Apply guardrails.
          4. NATO vocal challenge (Layer 3) for non-exempt tools.
          5. Run handler in thread-pool executor (Layer 4 — no event-loop blocking).
        """
        loop = asyncio.get_running_loop()
        self._loop = loop          # expose to sync handlers for fire-and-forget broadcasts
        tool_input = _strip_override(tool_name, tool_input)
        handler = getattr(self, f"_tool_{tool_name}", None)

        if handler is None:
            self._audit.log_action(
                tool_name, reasoning, "unknown", "error", "Tool no implementada"
            )
            return {"error": f"Tool '{tool_name}' no implementada."}

        preflight = self._preflight_validate(tool_name, tool_input)
        if preflight is not None:
            self._audit.log_action(
                tool_name, reasoning, "blocked:preflight", "blocked",
                preflight.get("error", "")[:200],
            )
            return preflight

        guardrail_block = self._validate_guardrails(tool_name, tool_input)
        if guardrail_block:
            self._audit.log_action(
                tool_name, reasoning, "blocked:guardrail", "blocked",
                guardrail_block.get("error", "")[:200],
            )
            return guardrail_block

        # V62.0 Phase 7 — risk-class taxonomy drives the actual gating decision.
        # classify_tool() is built to match _ALWAYS_HITL_TOOLS/_HITL_EXEMPT_TOOLS
        # exactly for every known tool (verified at import time above); this is
        # not a parallel decorative check, it IS the check now.
        risk_class = classify_tool(tool_name)
        if requires_trusted_lab(risk_class) and not _trusted_lab_enabled():
            self._audit.log_action(
                tool_name, reasoning, "blocked:lab_only", "blocked",
                "LAB_ONLY tool refused — JARVIS_TRUSTED_LAB is not enabled",
            )
            return {
                "error": (
                    f"Tool '{tool_name}' es LAB_ONLY — requiere "
                    "JARVIS_TRUSTED_LAB=true (y aun así, autorización HITL)."
                )
            }

        auth_audit = "hitl_exempt"
        must_challenge = requires_hitl(risk_class)
        if must_challenge:
            preview = str(tool_input)
            if len(preview) > 200:
                preview = preview[:200] + "…"
            from core.aura_events import ToolAuthPendingEvent
            await _aura_broadcast(ToolAuthPendingEvent(
                tool=tool_name,
                risk=risk_class.value,
                preview=preview,
                rollback_hint=rollback_hint(risk_class, tool_name),
            ).to_dict())
            granted, auth_audit = await self._challenge(tool_name, preview)
            if not granted:
                logger.warning(f"HITL: '{tool_name}' denegada.")
                self._audit.log_action(
                    tool_name, reasoning, auth_audit, "blocked", "Ejecución cancelada"
                )
                return {"error": "Ejecución cancelada por el usuario."}

        # Broadcast: tool is about to execute
        await _aura_broadcast({
            "type": "tool_invoked",
            "tool": tool_name,
            "auth_audit": auth_audit,
            "timestamp": datetime.now(timezone.utc).isoformat(),
        })

        try:
            # Layer 4: run synchronous tool handler in thread pool
            result = await loop.run_in_executor(None, lambda: handler(**tool_input))
            status = "error" if isinstance(result, dict) and "error" in result else "success"
            output_summary = json.dumps(result, ensure_ascii=False, default=str)[:200]
            result = self._check_pii_output(result)
            self._audit.log_action(tool_name, reasoning, auth_audit, status, output_summary)

            # Broadcast: tool result
            await _aura_broadcast({
                "type": "tool_result",
                "tool": tool_name,
                "status": status,
                "auth_audit": auth_audit,
                "timestamp": datetime.now(timezone.utc).isoformat(),
                "summary": output_summary[:150],
            })

            return result
        except Exception as e:
            logger.error(f"Error en tool '{tool_name}': {e}")
            self._audit.log_action(tool_name, reasoning, auth_audit, "error", str(e)[:200])
            await _aura_broadcast({
                "type": "error",
                "tool": tool_name,
                "message": str(e)[:200],
                "timestamp": datetime.now(timezone.utc).isoformat(),
            })
            return {"error": str(e)}

    async def aexecute_mcp(
        self, tool_name: str, tool_input: dict, call_fn, reasoning: str = ""
    ) -> Any:
        """
        Security gate for MCP-bridge tools — the counterpart to aexecute() for
        tools implemented by an external MCP server rather than a local
        _tool_* handler.

        *call_fn* is an async callable(tool_name, tool_input) that performs the
        actual RPC to the MCP server; it is NEVER invoked before the allowlist,
        filename-traversal, and HITL checks below all pass. MCP tools are
        foreign, unaudited code — allowlist-not-denylist. V62.0 Phase 7: HITL
        policy now comes from the SAME risk-class taxonomy aexecute() uses
        (core.risk_classes) — one gateway for both local and MCP dispatch.
        Both currently-allowlisted MCP tools classify REVERSIBLE, which still
        requires HITL unconditionally (no exempt tier for MCP), so this is not
        a behavior change; an unclassified future MCP tool defaults to
        HIGH_IMPACT (fail-closed) rather than silently skipping the challenge.
        """
        tool_input = _strip_override(tool_name, tool_input)

        if tool_name not in MCP_TOOL_ALLOWLIST:
            self._audit.log_action(
                tool_name, reasoning, "blocked:mcp_not_allowlisted", "blocked",
                "MCP tool no está en la allowlist explícita",
            )
            return {"error": f"Tool MCP '{tool_name}' no está permitida."}

        for key in _MCP_FILENAME_ARG_KEYS & tool_input.keys():
            err = _validate_mcp_filename(str(tool_input[key]))
            if err:
                self._audit.log_action(
                    tool_name, reasoning, "blocked:mcp_path_traversal", "blocked", err[:200],
                )
                return {"error": err}

        risk_class = classify_tool(tool_name)
        if requires_trusted_lab(risk_class) and not _trusted_lab_enabled():
            self._audit.log_action(
                tool_name, reasoning, "blocked:lab_only", "blocked",
                "LAB_ONLY MCP tool refused — JARVIS_TRUSTED_LAB is not enabled",
            )
            return {
                "error": (
                    f"Tool MCP '{tool_name}' es LAB_ONLY — requiere "
                    "JARVIS_TRUSTED_LAB=true (y aun así, autorización HITL)."
                )
            }

        preview = str(tool_input)
        if len(preview) > 200:
            preview = preview[:200] + "…"
        auth_audit = "hitl_exempt"
        if requires_hitl(risk_class):
            from core.aura_events import ToolAuthPendingEvent
            await _aura_broadcast(ToolAuthPendingEvent(
                tool=f"mcp:{tool_name}",
                risk=risk_class.value,
                preview=preview,
                rollback_hint=rollback_hint(risk_class, tool_name),
            ).to_dict())
            granted, auth_audit = await self._challenge(f"mcp:{tool_name}", preview)
            if not granted:
                logger.warning(f"HITL: MCP tool '{tool_name}' denegada.")
                self._audit.log_action(
                    tool_name, reasoning, auth_audit, "blocked", "Ejecución MCP cancelada por el usuario."
                )
                return {"error": "Ejecución cancelada por el usuario."}

        await _aura_broadcast({
            "type": "tool_invoked",
            "tool": f"mcp:{tool_name}",
            "auth_audit": auth_audit,
            "timestamp": datetime.now(timezone.utc).isoformat(),
        })

        try:
            result = await call_fn(tool_name, tool_input)
        except Exception as e:
            logger.error(f"Error en tool MCP '{tool_name}': {e}")
            self._audit.log_action(tool_name, reasoning, "mcp", "error", str(e)[:200])
            await _aura_broadcast({
                "type": "error",
                "tool": f"mcp:{tool_name}",
                "message": str(e)[:200],
                "timestamp": datetime.now(timezone.utc).isoformat(),
            })
            return {"error": f"MCP error en '{tool_name}': {e}"}

        status = "error" if isinstance(result, dict) and "error" in result else "success"
        result = self._check_pii_output(result)
        output_summary = json.dumps(result, ensure_ascii=False, default=str)[:200]
        self._audit.log_action(tool_name, reasoning, auth_audit, status, output_summary)

        await _aura_broadcast({
            "type": "tool_result",
            "tool": f"mcp:{tool_name}",
            "status": status,
            "auth_audit": auth_audit,
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "summary": output_summary[:150],
        })

        return result

    def _validate_guardrails(self, tool_name: str, tool_input: dict) -> dict | None:
        """Returns an error dict if the action violates a security guardrail, else None.

        SECURITY (V60.0): the destructive-pattern override is NOT taken from
        ``tool_input`` anymore. A model-generated tool argument such as
        ``FORCE_OVERRIDE=true`` can no longer disable these guardrails — the
        only legitimate override is operator-set trusted-lab mode, sourced
        exclusively from the environment / .env (see ``_trusted_lab_enabled``).
        """
        if "FORCE_OVERRIDE" in tool_input:
            # An override key reached the guardrail despite being stripped at the
            # gate — treat as an injection attempt, log loudly, and ignore it.
            logger.warning(
                f"Guardrail: ignored model-supplied FORCE_OVERRIDE on tool='{tool_name}' "
                "(LLM-controlled guardrail bypass is disabled)."
            )

        lab_override = _trusted_lab_enabled()

        combined = " ".join(
            str(v) for k, v in tool_input.items() if k != "FORCE_OVERRIDE"
        )

        if _GUARDRAIL_ROOT_DELETE_RE.search(combined):
            if lab_override:
                logger.warning(
                    f"Guardrail: root-delete pattern ALLOWED under trusted-lab mode — "
                    f"tool='{tool_name}' input={combined[:80]!r}"
                )
                return None
            logger.warning(
                f"Guardrail: eliminación de raíz bloqueada — tool='{tool_name}' "
                f"input={combined[:80]!r}"
            )
            return {
                "error": (
                    "GUARDRAIL: operación bloqueada — intento de eliminar un directorio raíz "
                    "detectado. Habilita JARVIS_TRUSTED_LAB=true en .env para anular "
                    "(solo en un laboratorio aislado y autorizado)."
                )
            }

        if _GUARDRAIL_SYSTEM_WRITE_RE.search(combined):
            if lab_override:
                logger.warning(
                    f"Guardrail: system-write pattern ALLOWED under trusted-lab mode — "
                    f"tool='{tool_name}' input={combined[:80]!r}"
                )
                return None
            logger.warning(
                f"Guardrail: escritura en ruta del sistema bloqueada — tool='{tool_name}' "
                f"input={combined[:80]!r}"
            )
            return {
                "error": (
                    "GUARDRAIL: operación bloqueada — modificación de C:\\Windows o System32 "
                    "no permitida. Habilita JARVIS_TRUSTED_LAB=true en .env para anular "
                    "(solo en un laboratorio aislado y autorizado)."
                )
            }

        return None

    def _check_pii_output(self, result: Any) -> Any:
        """Inject a PII warning key if the result contains cleartext credentials."""
        if not isinstance(result, dict):
            return result
        result_str = json.dumps(result, ensure_ascii=False, default=str)
        if _PII_OUTPUT_RE.search(result_str):
            logger.warning("PII detectado en output de tool — advertencia inyectada.")
            result["_pii_warning"] = (
                "ADVERTENCIA DE SEGURIDAD: Este resultado puede contener credenciales "
                "en texto claro (contraseñas, API keys o tokens). "
                "Revisa antes de mostrar o registrar en cualquier sistema externo."
            )
        return result

    def _cleanup_web_server(self) -> None:
        if self._active_web_server is not None:
            self._active_web_server.shutdown()
            self._active_web_server.server_close()
            if self._active_web_thread is not None:
                self._active_web_thread.join(timeout=3.0)
            self._active_web_server = None
            self._active_web_thread = None

    # ── Tiempo / Clima ────────────────────────────────────────────────────────

    def _tool_get_datetime(self) -> dict:
        now = datetime.now()
        return {
            "date": now.strftime("%A, %d de %B de %Y"),
            "time": now.strftime("%H:%M:%S"),
            "weekday": now.strftime("%A"),
        }

    def _tool_get_weather(self, city: str) -> dict:
        try:
            resp = requests.get(
                f"https://wttr.in/{requests.utils.quote(city)}?format=j1&lang=es",
                timeout=8,
            )
            resp.raise_for_status()
            data = resp.json()
            current = data["current_condition"][0]
            return {
                "city": city,
                "temp_c": current["temp_C"],
                "feels_like_c": current["FeelsLikeC"],
                "humidity": current["humidity"],
                "description": current["weatherDesc"][0]["value"],
                "wind_kmph": current["windspeedKmph"],
            }
        except Exception as e:
            return {"error": str(e)}

    # ── Lectura de archivos ───────────────────────────────────────────────────

    def _tool_read_file(self, path: str, max_chars: int = 8000) -> dict:
        """[SANDBOXED] Lee archivos solo dentro de Downloads, Documents o el proyecto."""
        p = Path(path).expanduser().resolve()
        allowed_dirs = [
            Path.home() / "Downloads",
            Path.home() / "Documents",
            Path.cwd(),
        ]
        is_allowed = any(p.is_relative_to(a.resolve()) for a in allowed_dirs)
        if not is_allowed:
            logger.warning(f"Intento de lectura bloqueado: {p}")
            return {"error": "Seguridad: No tengo permiso para leer fuera de Downloads o Documents."}
        if not p.exists():
            return {"error": f"Archivo no encontrado: {path}"}

        ext = p.suffix.lower()
        try:
            if ext == ".pdf":
                content = self._read_pdf(p)
            elif ext in (".docx", ".doc"):
                content = self._read_docx(p)
            elif ext in (".xlsx", ".xls"):
                content = self._read_xlsx(p)
            elif ext in (".pptx", ".ppt"):
                content = self._read_pptx(p)
            elif ext in (
                ".txt", ".md", ".py", ".js", ".sh", ".yaml", ".yml",
                ".json", ".xml", ".html", ".css", ".log", ".c",
                ".cpp", ".h", ".java", ".go", ".rs", ".env", ".csv",
            ):
                content = p.read_text(encoding="utf-8", errors="ignore")
            elif ext == ".rtf":
                content = self._read_rtf(p)
            elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"):
                content = self._read_image_ocr(p)
            else:
                return {"error": f"Extensión '{ext}' no soportada."}

            if len(content) > max_chars:
                content = content[:max_chars] + f"\n\n[...truncado a {max_chars} chars]"

            return {
                "file": str(p.name),
                "extension": ext,
                "size_kb": round(p.stat().st_size / 1024, 2),
                "content": content,
                "chars": len(content),
            }
        except Exception as e:
            return {"error": f"Error leyendo {path}: {e}"}

    def _read_pdf(self, path: Path) -> str:
        import pdfplumber
        parts = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                if text:
                    parts.append(f"[Pag {i}]\n{text}")
                for table in page.extract_tables():
                    rows = [" | ".join(str(c) for c in row if c) for row in table if row]
                    parts.append("[Tabla]\n" + "\n".join(rows))
        return "\n\n".join(parts)

    def _read_docx(self, path: Path) -> str:
        from docx import Document
        doc = Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text for c in row.cells))
        return "\n".join(parts)

    def _read_xlsx(self, path: Path) -> str:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for name in wb.sheetnames:
            parts.append(f"[Hoja: {name}]")
            for row in wb[name].iter_rows(values_only=True):
                if any(c is not None for c in row):
                    parts.append(" | ".join(str(c) if c is not None else "" for c in row))
        return "\n".join(parts)

    def _read_pptx(self, path: Path) -> str:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                s.text.strip()
                for s in slide.shapes
                if hasattr(s, "text") and s.text.strip()
            ]
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts)

    def _read_rtf(self, path: Path) -> str:
        from striprtf.striprtf import rtf_to_text
        return rtf_to_text(path.read_text(errors="ignore"))

    def _read_image_ocr(self, path: Path) -> str:
        try:
            import pytesseract
            from PIL import Image
            return pytesseract.image_to_string(Image.open(str(path)), lang="spa+eng")
        except Exception:
            pass
        try:
            import easyocr
            reader = easyocr.Reader(["es", "en"], gpu=False)
            return " ".join(reader.readtext(str(path), detail=0))
        except Exception as e:
            return f"OCR no disponible: {e}"

    def _tool_list_directory(self, path: str = ".", pattern: str = "*") -> dict:
        p = Path(path).expanduser()
        if not p.exists():
            return {"error": f"Directorio no encontrado: {path}"}
        files = [
            {
                "name": item.name,
                "type": "dir" if item.is_dir() else "file",
                "ext": item.suffix.lower(),
                "size_kb": round(item.stat().st_size / 1024, 2) if item.is_file() else 0,
            }
            for item in sorted(p.glob(pattern))
        ]
        return {"path": str(p.resolve()), "items": files, "count": len(files)}

    # ── Lectura Universal y SAST ──────────────────────────────────────────────

    def _tool_leer_archivo_universal(self, filepath: str) -> dict:
        """Lee archivos multiformato con truncamiento estricto de 4000 chars (Single Channel VRAM)."""
        _MAX_CHARS = 4000
        p = Path(filepath).expanduser().resolve()
        if not p.exists():
            return {"error": f"Archivo no encontrado: {filepath}"}

        ext = p.suffix.lower()
        try:
            if ext == ".pdf":
                import pdfplumber
                parts = []
                with pdfplumber.open(p) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text() or ""
                        if text:
                            parts.append(text)
                content = "\n\n".join(parts)
            elif ext in (".docx", ".doc"):
                from docx import Document
                doc = Document(str(p))
                content = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
            elif ext in (".csv", ".xlsx", ".xls"):
                import pandas as pd
                if ext == ".csv":
                    df = pd.read_csv(p, nrows=50)
                else:
                    df = pd.read_excel(p, nrows=50)
                content = df.to_string(index=False)
            else:
                content = p.read_text(encoding="utf-8", errors="ignore")
        except Exception as e:
            return {"error": f"Error leyendo {filepath}: {e}"}

        result: dict = {
            "file": p.name,
            "extension": ext,
            "chars": min(len(content), _MAX_CHARS),
            "content": content[:_MAX_CHARS],
        }
        if len(content) > _MAX_CHARS:
            result["aviso"] = (
                f"Texto truncado a {_MAX_CHARS} caracteres para proteger la memoria (VRAM)."
            )
        return result

    def _tool_analizar_codigo_sast(self, filepath: str) -> dict:
        """Análisis estático ligero (regex-based SAST) sin dependencias externas."""
        _MAX_FINDINGS = 15
        p = Path(filepath).expanduser().resolve()
        if not p.exists():
            return {"error": f"Archivo no encontrado: {filepath}"}

        findings: list[dict] = []
        possibly_more = False
        try:
            with p.open(encoding="utf-8", errors="ignore") as f:
                for lineno, line in enumerate(f, 1):
                    for pattern, label in _SAST_PATTERNS:
                        if pattern.search(line):
                            findings.append({
                                "linea": lineno,
                                "tipo": label,
                                "codigo": line.rstrip(),
                            })
                            break
                    if len(findings) >= _MAX_FINDINGS:
                        possibly_more = True
                        break
        except Exception as e:
            return {"error": f"Error analizando {filepath}: {e}"}

        if not findings:
            return {"status": "limpio", "archivo": p.name, "hallazgos": 0}

        result: dict = {"archivo": p.name, "hallazgos": len(findings), "resultados": findings}
        if possibly_more:
            result["aviso"] = "Límite de 15 hallazgos alcanzado; puede haber más en el resto del archivo."
        return result

    # ── Búsqueda Web ──────────────────────────────────────────────────────────

    def _tool_web_search(self, query: str, max_results: int = 5) -> dict:
        try:
            from duckduckgo_search import DDGS
            with DDGS() as ddg:
                results = list(ddg.text(query, max_results=max_results))
            return {
                "query": query,
                "results": [
                    {"title": r["title"], "url": r["href"], "snippet": r["body"]}
                    for r in results
                ],
            }
        except Exception as e:
            return {"error": str(e)}

    def _tool_fetch_webpage(self, url: str, max_chars: int = 5000) -> dict:
        try:
            from bs4 import BeautifulSoup
            resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=10)
            resp.raise_for_status()
            soup = BeautifulSoup(resp.text, "lxml")
            for tag in soup(["script", "style", "nav", "footer"]):
                tag.decompose()
            text = soup.get_text(separator="\n", strip=True)
            return {"url": url, "content": text[:max_chars]}
        except Exception as e:
            return {"error": str(e)}

    # ── Sistema ───────────────────────────────────────────────────────────────

    def _tool_system_info(self) -> dict:
        import psutil
        cpu = psutil.cpu_percent(interval=1)
        ram = psutil.virtual_memory()
        disk = psutil.disk_usage("/")
        uptime = datetime.now() - datetime.fromtimestamp(psutil.boot_time())
        return {
            "cpu_percent": cpu,
            "ram_total_gb": round(ram.total / 1e9, 2),
            "ram_used_gb": round(ram.used / 1e9, 2),
            "ram_percent": ram.percent,
            "disk_total_gb": round(disk.total / 1e9, 2),
            "disk_percent": disk.percent,
            "uptime_hours": round(uptime.total_seconds() / 3600, 1),
            "os": platform.system(),
            "hostname": platform.node(),
        }

    def _tool_list_processes(self, filter_name: str = "") -> dict:
        import psutil
        procs = []
        for p in psutil.process_iter(["pid", "name", "cpu_percent", "memory_percent"]):
            try:
                if filter_name.lower() in p.info["name"].lower():
                    procs.append(p.info)
            except Exception:
                pass
        procs.sort(key=lambda x: x.get("cpu_percent", 0), reverse=True)
        return {"processes": procs[:30]}

    def _tool_kill_process(self, name: str) -> dict:
        import psutil
        killed = []
        for p in psutil.process_iter(["pid", "name"]):
            try:
                if name.lower() in p.info["name"].lower():
                    p.kill()
                    killed.append(p.info["name"])
            except Exception:
                pass
        return {"killed": killed}

    def _tool_run_shell_command(self, command: str) -> dict:
        """
        [ALLOWLIST + Layer1 flags + Layer2 canonicalization + shell=False]

        Authorization is handled by the outer async _challenge() NATO gate.
        This method validates the command and executes it — no blocking input().
        """
        is_valid, error_msg, argv = _validate_command(command)
        if not is_valid:
            logger.warning(f"Comando bloqueado: {command!r} — {error_msg}")
            self._audit.log_action(
                "run_shell_command", "", "blocked:validation", "blocked", error_msg[:200],
                command=command,
            )
            # Static forensic triage — zero execution, analysis only
            try:
                from core.triage import analyze_neutralized, write_manifest

                # Layer 0: YARA scan — compiled rules are lru_cache'd (one-time cost)
                yara_hits: list[dict] = []
                try:
                    from core.yara_analyzer import _compile_rules
                    _rules = _compile_rules()
                    if _rules:
                        yara_hits = [
                            {"rule": m.rule, "namespace": m.namespace, "tags": list(m.tags)}
                            for m in _rules.match(data=command.encode())
                        ]
                        if yara_hits:
                            logger.warning(
                                f"YARA: {len(yara_hits)} hit(s) — "
                                f"{[h['rule'] for h in yara_hits]}"
                            )
                except Exception as _ye:
                    logger.warning(f"YARA scan error: {_ye}")

                triage_result  = analyze_neutralized(command, error_msg, yara_hits=yara_hits)
                manifest_path  = write_manifest(triage_result, command)

                # Episodic memory — store from executor thread via threadsafe bridge
                _loop = getattr(self, "_loop", None)
                if _loop and not _loop.is_closed():
                    try:
                        from core.episodic_memory import store_episode as _store_ep
                        import json as _json_ep
                        asyncio.run_coroutine_threadsafe(
                            _store_ep(
                                _json_ep.dumps(triage_result, ensure_ascii=False, default=str),
                                "triage",
                                severity="HIGH",
                                mitre_tags=triage_result.get("mitre_match", []),
                            ),
                            _loop,
                        )
                    except Exception:
                        pass

                with open("tactic_audit.jsonl", "a", encoding="utf-8") as _af:
                    _af.write(json.dumps({
                        "status":    "neutralized",
                        "command":   command,
                        "triage":    triage_result,
                        "manifest":  manifest_path.name,
                        "timestamp": datetime.now(timezone.utc).isoformat(),
                    }, ensure_ascii=False) + "\n")
                _loop = getattr(self, "_loop", None)
                if _loop and not _loop.is_closed():
                    asyncio.run_coroutine_threadsafe(
                        _aura_broadcast({
                            "type":          "neutralized",
                            "command":       command,
                            "triage":        triage_result,
                            "manifest_file": manifest_path.name,
                            "timestamp":     datetime.now(timezone.utc).isoformat(),
                        }),
                        _loop,
                    )

                # SOAR: isolate public IPs that clear the entropy + MITRE AND-gate
                try:
                    from core.mitigation import should_isolate, isolate_ip
                    target_ips = should_isolate(triage_result)
                    if target_ips and _loop and not _loop.is_closed():
                        for _ip in target_ips:
                            logger.warning(f"SOAR: scheduling firewall isolation for {_ip}")

                            async def _run_isolation(_i=_ip):
                                asyncio.create_task(
                                    asyncio.shield(isolate_ip(_i, _aura_broadcast))
                                )

                            asyncio.run_coroutine_threadsafe(_run_isolation(), _loop)
                except Exception as _soar_exc:
                    logger.warning(f"SOAR pipeline error: {_soar_exc}")

                logger.info(f"Triage manifest written: {manifest_path.name}")
            except Exception as _triage_exc:
                logger.warning(f"Triage pipeline error: {_triage_exc}")
            return {"error": f"Comando bloqueado por política de seguridad: {error_msg}"}

        # Show the canonicalized argv for operator transparency. V62.0 Phase 7:
        # binary_risk_class() is informational only here (audit trail) — it
        # does not gate anything; run_shell_command itself is already
        # unconditionally HIGH_IMPACT (see core/risk_classes.py), and
        # _validate_command's allowlist + core/trust_engine.py's dynamic
        # trust floor remain the sole authorities over which binaries run.
        binary = Path(argv[0]).name.lower().removesuffix(".exe") if argv else ""
        print(f"\n    [EXEC] argv={argv} risk={binary_risk_class(binary).value}")
        sys.stdout.flush()

        try:
            result = subprocess.run(
                argv,
                shell=False,
                capture_output=True,
                text=True,
                timeout=30,
            )
            return {
                "stdout": result.stdout.strip(),
                "stderr": result.stderr.strip(),
                "returncode": result.returncode,
            }
        except FileNotFoundError:
            return {"error": f"Ejecutable '{argv[0]}' no encontrado en PATH."}
        except subprocess.TimeoutExpired:
            return {"error": "Timeout: el comando tardó más de 30 segundos."}
        except Exception as e:
            return {"error": str(e)}

    # ── Aplicaciones ──────────────────────────────────────────────────────────

    def _tool_open_application(self, app: str) -> dict:
        OS = platform.system()
        # Todos los valores son listas pre-definidas — el input del usuario
        # solo se usa como clave de búsqueda, nunca se interpola en el comando.
        APP_MAP: dict[str, dict[str, list[str]]] = {
            "packet tracer": {
                "Windows": [r"C:\Program Files\Cisco\Cisco Packet Tracer 8.2\PacketTracer.exe"],
                "Linux": ["packettracer"],
            },
            "wireshark":  {"Windows": ["wireshark"],    "Linux": ["wireshark"]},
            "burpsuite":  {"Windows": ["burpsuite"],    "Linux": ["burpsuite"]},
            "vscode":     {"Windows": ["code"],         "Linux": ["code"]},
            "chrome":     {"Windows": ["chrome"],       "Linux": ["google-chrome"]},
            "firefox":    {"Windows": ["firefox"],      "Linux": ["firefox"]},
            "terminal":   {"Windows": ["cmd.exe"],      "Linux": ["x-terminal-emulator"]},
            "calculator": {"Windows": ["calc.exe"],     "Linux": ["gnome-calculator"]},
            "excel":      {"Windows": ["excel.exe"],    "Linux": ["libreoffice", "--calc"]},
            "word":       {"Windows": ["winword.exe"],  "Linux": ["libreoffice", "--writer"]},
        }

        key = app.lower().strip()
        if len(key) > 64:
            return {"error": "Nombre de aplicación demasiado largo."}

        cmd_vector: list[str] | None = None
        for k, v in APP_MAP.items():
            if key in k or k in key:
                cmd_vector = v.get(OS, v.get("Linux"))
                break

        if cmd_vector is None:
            # Fuera del mapa: solo permitir ejecutables sin metacaracteres ni espacios
            if not re.match(r'^[a-zA-Z0-9._-]+$', key):
                return {"error": "Nombre de aplicación contiene caracteres no permitidos."}
            cmd_vector = [key]

        try:
            subprocess.Popen(
                cmd_vector,
                shell=False,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
            return {"status": "opened", "app": app}
        except FileNotFoundError:
            return {"error": f"Aplicación '{cmd_vector[0]}' no encontrada en PATH."}
        except Exception as e:
            return {"error": str(e)}

    def _tool_open_software(self, app_name: str) -> dict:
        """
        [HITL] Lanza una aplicación usando un mapa pre-definido de ejecutables.
        El input del usuario solo se usa como clave de búsqueda — nunca se interpola.
        Fallback seguro: os.startfile() para aliases registrados en Windows.
        shell=False siempre.
        """
        key = app_name.lower().strip()
        if len(key) > 64:
            return {"error": "Nombre de aplicación demasiado largo."}

        # Exact match primero, luego partial match — el resultado es siempre un valor pre-definido
        executable = _APP_SOFTWARE_MAP.get(key)
        if executable is None:
            for k, v in _APP_SOFTWARE_MAP.items():
                if key in k or k in key:
                    executable = v
                    break

        # Si no está en el mapa, validar que no tenga metacaracteres antes de usarlo
        if executable is None:
            if not re.match(r'^[a-zA-Z0-9._-]+$', key):
                return {"error": "Nombre de aplicación contiene caracteres no permitidos."}
            executable = key

        try:
            subprocess.Popen(
                [executable],
                shell=False,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
            return {"status": "launched", "app": app_name, "executable": executable}
        except FileNotFoundError:
            pass
        except Exception as e:
            return {"error": str(e)}

        # Fallback seguro para Windows: ShellExecute vía os.startfile()
        if platform.system() == "Windows":
            try:
                os.startfile(executable)
                return {"status": "launched_via_startfile", "app": app_name, "executable": executable}
            except Exception as e:
                return {"error": f"Aplicación '{executable}' no encontrada. {e}"}

        return {"error": f"Aplicación '{executable}' no encontrada en el PATH."}

    def _tool_create_document(self, doc_type: str, title: str, content: str | list) -> dict:
        """
        [HITL] Genera un archivo .docx o .pptx en ~/Downloads.
        Para pptx, fragmenta el contenido en slides de máximo 5 bullets para
        evitar overflow de texto. Retorna la ruta absoluta del archivo creado.
        """
        slug = re.sub(r'[^\w\s-]', '', title).strip()
        slug = re.sub(r'[\s-]+', '_', slug)[:60] or "documento"

        out_dir = Path.home() / "Downloads"
        out_dir.mkdir(parents=True, exist_ok=True)

        lines: list[str] = []
        if isinstance(content, str):
            lines = [ln for ln in content.split('\n') if ln.strip()]
        elif isinstance(content, list):
            lines = [str(item) for item in content if str(item).strip()]

        if doc_type == "docx":
            from docx import Document
            doc = Document()
            doc.add_heading(title, level=1)
            for line in lines:
                doc.add_paragraph(line)
            out_path = out_dir / f"{slug}.docx"
            doc.save(str(out_path))
            return {"status": "created", "path": str(out_path.resolve()), "format": "docx"}

        elif doc_type == "pptx":
            from pptx import Presentation
            prs = Presentation()

            # Slide de título
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = title
            try:
                title_slide.placeholders[1].text = ""
            except (KeyError, IndexError):
                pass

            # Fragmentar contenido: máx 5 bullets por slide
            CHUNK = 5
            chunks = [lines[i:i + CHUNK] for i in range(0, len(lines), CHUNK)] if lines else [[]]
            content_layout = prs.slide_layouts[1]
            for idx, chunk in enumerate(chunks, 1):
                cslide = prs.slides.add_slide(content_layout)
                cslide.shapes.title.text = title if idx == 1 else f"{title} ({idx})"
                try:
                    tf = cslide.placeholders[1].text_frame
                    tf.clear()
                    for j, bullet in enumerate(chunk):
                        if j == 0:
                            tf.paragraphs[0].text = bullet
                        else:
                            tf.add_paragraph().text = bullet
                except (KeyError, IndexError):
                    pass

            out_path = out_dir / f"{slug}.pptx"
            prs.save(str(out_path))
            return {"status": "created", "path": str(out_path.resolve()), "format": "pptx"}

        else:
            return {"error": f"doc_type '{doc_type}' no soportado. Usa 'docx' o 'pptx'."}

    def _tool_desplegar_webapp(self, html_code: str) -> dict:
        self._cleanup_web_server()

        if re.search(r'<head\b', html_code, re.IGNORECASE):
            html_code = re.sub(
                r'(<head[^>]*>)',
                rf'\1\n  {_CSP_META}',
                html_code,
                flags=re.IGNORECASE,
            )
        else:
            html_code = _CSP_META + "\n" + html_code

        temp_dir = tempfile.mkdtemp(prefix="jarvis_ui_")
        (Path(temp_dir) / "index.html").write_text(html_code, encoding="utf-8")

        class JailedHTTPRequestHandler(http.server.SimpleHTTPRequestHandler):
            def __init__(self, *args, **kwargs):
                super().__init__(*args, directory=temp_dir, **kwargs)

            def log_message(self, format, *args):
                pass

        httpd = socketserver.TCPServer(("", 0), JailedHTTPRequestHandler)
        port = httpd.server_address[1]

        thread = threading.Thread(target=httpd.serve_forever, daemon=True)
        thread.start()

        self._active_web_server = httpd
        self._active_web_thread = thread

        url = f"http://localhost:{port}/index.html"
        webbrowser.open(url)

        return {"url": url, "port": port, "temp_dir": temp_dir}

    def _tool_take_screenshot(
        self,
        save_path: str = "",
        analyze: bool = False,
        analizar_topologia: bool = False,
    ) -> dict:
        if not self.consent.screen:
            return self._consent_error("screen capture", "enable screen access")

        # [SANDBOXED] Contain the caller-supplied save_path to the same allowed
        # roots as read_file/write_file. Absent a path, default to a timestamped
        # PNG under Downloads (an allowed root) — never the bare home dir, which
        # is outside containment. Validate BEFORE importing/invoking pyautogui so
        # a rejected path never captures the screen.
        if not save_path:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            target = Path.home() / "Downloads" / f"jarvis_{ts}.png"
        else:
            resolved = _resolve_within_allowed(save_path)
            if resolved is None:
                logger.warning(f"Screenshot save blocked (outside allowed dirs): {save_path!r}")
                return {"error": "Seguridad: solo puedo guardar capturas en Downloads, Documents o el proyecto."}
            target = resolved

        import pyautogui

        target.parent.mkdir(parents=True, exist_ok=True)
        screenshot = pyautogui.screenshot()
        screenshot.save(str(target))
        result: dict = {"saved": str(target)}
        if analyze:
            result["ocr"] = self._read_image_ocr(target)[:2000]
        if analizar_topologia:
            result["topology_analysis"] = self._analyze_topology_vlm(target)
        return result

    def _analyze_topology_vlm(self, image_path: Path) -> str:
        """Envía la imagen al VLM llava de Ollama para análisis de diagramas de red."""
        import base64

        try:
            with open(image_path, "rb") as f:
                image_b64 = base64.b64encode(f.read()).decode("utf-8")

            payload = {
                "model": "llava",
                "prompt": (
                    "Analiza este diagrama de red en detalle. "
                    "Identifica y lista: routers (con IPs si son visibles), switches, "
                    "hosts/PCs, conexiones entre dispositivos, VLANs, y la topología general. "
                    "Si hay configuraciones de red visibles (máscaras, gateways, protocolos), "
                    "inclúyelas. Responde en español de forma concisa y estructurada."
                ),
                "images": [image_b64],
                "stream": False,
            }
            resp = requests.post(
                "http://localhost:11434/api/generate",
                json=payload,
                timeout=90,
            )
            resp.raise_for_status()
            return resp.json().get("response", "Sin análisis disponible.")
        except Exception as e:
            return f"Error en análisis VLM: {e}"

    def _tool_escanear_pantalla(self) -> dict:
        """[OCR] Captura la pantalla y extrae el texto visible vía pytesseract (CPU-only)."""
        if not self.consent.screen:
            return self._consent_error("screen capture", "enable screen access")
        _MAX_CHARS = 3000
        try:
            from PIL import ImageGrab
        except ImportError:
            return {"error": "Pillow (PIL.ImageGrab) no instalado. pip install Pillow"}
        try:
            import pytesseract
        except ImportError:
            return {"error": "pytesseract no instalado. pip install pytesseract (requiere tesseract en el OS)"}

        try:
            img = ImageGrab.grab()
        except Exception as e:
            return {"error": f"No se pudo capturar la pantalla: {e}"}

        try:
            raw_text = pytesseract.image_to_string(img, lang="spa+eng")
        except pytesseract.TesseractNotFoundError:
            return {"error": "Binario tesseract no encontrado en PATH. Instala tesseract-ocr en el sistema."}
        except Exception as e:
            return {"error": f"OCR falló: {e}"}

        cleaned_lines = [
            "".join(c for c in line if c.isprintable()).strip()
            for line in raw_text.splitlines()
        ]
        cleaned = "\n".join(line for line in cleaned_lines if line)

        if not cleaned:
            return {"status": "ilegible", "aviso": "No se detectó texto en la pantalla."}

        truncated = len(cleaned) > _MAX_CHARS
        if truncated:
            cleaned = cleaned[:_MAX_CHARS]

        result: dict = {
            "status": "ok",
            "chars": len(cleaned),
            "texto": cleaned,
        }
        if truncated:
            result["aviso"] = f"Texto truncado a {_MAX_CHARS} caracteres para proteger el KV Cache."
        return result

    def _tool_press_hotkey(self, keys: list) -> dict:
        import pyautogui
        pyautogui.hotkey(*keys)
        return {"pressed": "+".join(keys)}

    def _tool_type_text(self, text: str) -> dict:
        import pyautogui
        pyautogui.write(text, interval=0.05)
        return {"typed": len(text)}

    # ── Clipboard ─────────────────────────────────────────────────────────────

    def _tool_get_clipboard(self) -> dict:
        if not self.consent.clipboard:
            return self._consent_error("clipboard access", "enable clipboard access")
        import pyperclip
        return {"clipboard": pyperclip.paste()}

    def _tool_set_clipboard(self, text: str) -> dict:
        if not self.consent.clipboard:
            return self._consent_error("clipboard access", "enable clipboard access")
        import pyperclip
        pyperclip.copy(text)
        return {"status": "copied", "length": len(text)}

    # ── Packet Tracer / Networking ────────────────────────────────────────────

    def _tool_packet_tracer_open(self, file_path: str = "") -> dict:
        OS = platform.system()
        candidates: dict[str, list[str]] = {
            "Windows": [
                r"C:\Program Files\Cisco\Cisco Packet Tracer 8.2\PacketTracer.exe",
                r"C:\Program Files\Cisco\Cisco Packet Tracer 8.1\PacketTracer.exe",
            ],
            "Linux": ["packettracer", "/opt/pt/bin/PacketTracer"],
        }
        pt_cmd: str | None = None
        for c in candidates.get(OS, candidates["Linux"]):
            if shutil.which(c) or Path(c).exists():
                pt_cmd = c
                break
        if not pt_cmd:
            return {"error": "Packet Tracer no encontrado en el PATH."}
        args = [pt_cmd] + ([file_path] if file_path else [])
        subprocess.Popen(args, shell=False)
        return {"status": "launched", "file": file_path or "nuevo proyecto"}

    def _tool_network_scan(self, target: str, scan_type: str = "-sS -sV") -> dict:
        """[VALIDATED] Target and scan_type validated before passing to python-nmap."""
        try:
            target = _validate_network_target(target)
        except ValueError as e:
            return {"error": str(e)}

        if _contains_shell_metacharacters(scan_type):
            return {"error": "scan_type inválido: contiene metacaracteres de shell prohibidos."}

        try:
            import nmap
        except ImportError:
            return {"error": "python-nmap not installed. Run: pip install python-nmap"}

        try:
            nm = nmap.PortScanner()
            nm.scan(hosts=target, arguments=scan_type)

            results: dict = {}
            for host in nm.all_hosts():
                host_data: dict = {
                    "state": nm[host].state(),
                    "protocols": {},
                }
                for proto in nm[host].all_protocols():
                    ports: dict = {}
                    for port in sorted(nm[host][proto].keys()):
                        port_info = nm[host][proto][port]
                        ports[port] = {
                            "state": port_info.get("state", ""),
                            "name": port_info.get("name", ""),
                            "product": port_info.get("product", ""),
                            "version": port_info.get("version", ""),
                        }
                    host_data["protocols"][proto] = ports
                results[host] = host_data

            return {"target": target, "scan_type": scan_type, "hosts": results}
        except nmap.PortScannerError as e:
            return {
                "error": (
                    "Nmap executable not found on host OS. "
                    f"Install it from https://nmap.org/download.html — Details: {e}"
                )
            }
        except Exception as e:
            return {"error": str(e)}

    def _tool_check_connectivity(self, host: str, port: int = 0) -> dict:
        """[VALIDATED] Host validado antes de usarlo en subprocess."""
        import socket

        if not _SAFE_HOST_RE.match(host):
            return {"error": "Host inválido. Solo IPs o hostnames válidos."}

        if port == 0:
            cmd_vector = (
                ["ping", "-c", "1", host]
                if platform.system() != "Windows"
                else ["ping", "-n", "1", host]
            )
            try:
                r = subprocess.run(
                    cmd_vector, shell=False, capture_output=True, text=True, timeout=10
                )
                return {"host": host, "reachable": r.returncode == 0}
            except Exception as e:
                return {"error": str(e)}

        try:
            socket.create_connection((host, port), timeout=5).close()
            return {"host": host, "port": port, "open": True}
        except Exception as e:
            return {"host": host, "port": port, "open": False, "error": str(e)}

    def _tool_estudiar_tema(self, url: str) -> dict:
        """Descarga una URL, limpia el HTML y vectoriza el contenido en ChromaDB."""
        import hashlib
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            return {"error": "beautifulsoup4 no instalado. Ejecuta: pip install beautifulsoup4 lxml"}

        try:
            resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=15)
            resp.raise_for_status()
            soup = BeautifulSoup(resp.text, "lxml")
            for tag in soup(["script", "style", "nav", "footer", "aside"]):
                tag.decompose()
            text = soup.get_text(separator="\n", strip=True)
        except Exception as e:
            return {"error": f"No se pudo descargar {url}: {e}"}

        if not text.strip():
            return {"error": "No se pudo extraer texto de la URL."}

        # V62.0 Phase 3 — web content is untrusted (core.memory_router.
        # is_untrusted_source models "web"/"url" as such); reject obvious
        # prompt-injection payloads at ingest time rather than relying solely
        # on the retrieval-time untrusted-tool-output banner (core/llm.py's
        # _UNTRUSTED_TOOL_SOURCES already labels consultar_base_conocimiento's
        # results, but that's advisory — this is a hard gate before storage).
        from core.feed_sanitizer import check_prompt_injection, SanitizationError
        try:
            check_prompt_injection(text, source=url)
        except SanitizationError:
            return {"error": f"Contenido de {url} rechazado: posible inyección de prompt detectada."}

        if not hasattr(self, "_memory"):
            try:
                from core.memory import VectorMemory
                self._memory = VectorMemory()
            except Exception as e:
                return {"error": f"No se pudo inicializar la memoria vectorial: {e}"}

        words = text.split()
        chunk_size, overlap = 400, 50
        chunks_indexed = 0
        i = 0
        while i < len(words):
            chunk = " ".join(words[i : i + chunk_size])
            chunk_id = hashlib.md5(f"{url}:{i}".encode()).hexdigest()
            self._memory.add(chunk, chunk_id, {"source": url, "chunk": chunks_indexed})
            chunks_indexed += 1
            i += chunk_size - overlap

        return {
            "url": url,
            "chars_extraidos": len(text),
            "chunks_indexados": chunks_indexed,
            "status": "Listo. Usa consultar_base_conocimiento para recuperar el contenido.",
        }

    def _tool_consultar_base_conocimiento(self, query: str) -> dict:
        """Vectoriza la pregunta y recupera los 3 fragmentos más relevantes de ChromaDB."""
        if not hasattr(self, "_memory"):
            try:
                from core.memory import VectorMemory
                self._memory = VectorMemory()
            except Exception as e:
                return {"error": f"No se pudo inicializar la memoria vectorial: {e}"}
        return self._memory.query(query)

    def _tool_whois_lookup(self, domain: str) -> dict:
        """[VALIDATED] Dominio validado contra regex antes de pasarlo a whois."""
        if not _SAFE_DOMAIN_RE.match(domain):
            return {
                "error": "Dominio inválido. Solo alfanuméricos, puntos y guiones (1-253 chars)."
            }
        try:
            r = subprocess.run(
                ["whois", domain], shell=False, capture_output=True, text=True, timeout=15
            )
            return {"domain": domain, "whois": r.stdout[:3000]}
        except FileNotFoundError:
            return {"error": "whois no encontrado. Instálalo con: apt install whois"}
        except Exception as e:
            return {"error": str(e)}

    # ── Knowledge Vault (RAG local) ───────────────────────────────────────────

    def _get_vault(self):
        if not hasattr(self, "_vault"):
            from core.knowledge import KnowledgeVault
            self._vault = KnowledgeVault()
        return self._vault

    def _tool_ingest_docs(self, folder_path: str = "") -> dict:
        """Index PDFs and TXTs from a local folder into the Knowledge Vault."""
        return self._get_vault().ingest_docs(folder_path)

    def _tool_query_knowledge(self, query: str) -> dict:
        """Retrieve the top-3 most relevant chunks from the Knowledge Vault."""
        text = self._get_vault().query_knowledge(query)
        return {"result": text}

    # ── System Health Monitor ─────────────────────────────────────────────────

    def _tool_get_system_status(self) -> dict:
        """CPU, RAM, battery and theoretical memory bandwidth saturation report."""
        cpu_pct = psutil.cpu_percent(interval=0.5)
        ram = psutil.virtual_memory()

        # Theoretical estimate: RAM occupancy weighted with CPU pressure.
        bw_saturation = round(min(ram.percent * 0.70 + cpu_pct * 0.30, 100.0), 1)

        report: dict = {
            "cpu_usage_pct": cpu_pct,
            "ram": {
                "total_gb": round(ram.total / 1e9, 2),
                "used_gb": round(ram.used / 1e9, 2),
                "available_gb": round(ram.available / 1e9, 2),
                "usage_pct": ram.percent,
            },
            "memory_bandwidth_saturation_estimate_pct": bw_saturation,
        }

        if ram.percent > 85:
            report["warning"] = (
                "RAM usage crítica (>85%). "
                "Considera cerrar Chrome, Blender u otras apps secundarias."
            )

        try:
            battery = psutil.sensors_battery()
            if battery:
                secsleft = battery.secsleft
                report["battery"] = {
                    "percent": battery.percent,
                    "plugged_in": battery.power_plugged,
                    "time_left_min": round(secsleft / 60, 1) if secsleft > 0 else "N/A",
                }
            else:
                report["battery"] = {"status": "No detectada (posiblemente desktop)"}
        except Exception:
            report["battery"] = {"status": "No disponible en este sistema"}

        governance_ok = self._audit.is_writable()
        report["governance_check"] = {
            "audit_logger_active": True,
            "audit_log_writable": governance_ok,
            "audit_log_path": str(self._audit.log_path),
            "status": "OPERATIONAL" if governance_ok else "DEGRADED — log file not writable",
        }

        return report

    def _tool_osint_lookup(self, domain: str) -> dict:
        """[HITL + VALIDATED] WHOIS + DNS recon via python-whois and dnspython."""
        # Strip protocol prefix and www
        domain = re.sub(r'^https?://', '', domain)
        domain = re.sub(r'^www\.', '', domain)
        domain = domain.split('/')[0].strip()

        if not _SAFE_DOMAIN_RE.match(domain):
            return {"error": "Dominio inválido. Solo alfanuméricos, puntos y guiones (1-253 chars)."}

        result: dict = {"domain": domain}

        try:
            w = whois.whois(domain)
            result["whois"] = {
                "registrar": w.registrar,
                "creation_date": str(w.creation_date),
                "expiration_date": str(w.expiration_date),
                "name_servers": list(w.name_servers) if w.name_servers else None,
            }
        except Exception as e:
            result["whois"] = {"error": str(e)}

        dns_records: dict = {}
        for record_type in ("A", "MX", "TXT"):
            try:
                answers = dns.resolver.resolve(domain, record_type)
                dns_records[record_type] = [str(r) for r in answers]
            except Exception as e:
                dns_records[record_type] = {"error": str(e)}
        result["dns"] = dns_records

        return result

    # ── V59.0 APEX — Power Tools ──────────────────────────────────────────────

    def _tool_write_file(self, path: str, content: str, mode: str = "w") -> dict:
        """[HITL] Write text content to a file in Downloads, Documents, or project dir."""
        p = Path(path).expanduser().resolve()
        allowed_dirs = [
            Path.home() / "Downloads",
            Path.home() / "Documents",
            Path.cwd(),
        ]
        if not any(p.is_relative_to(a.resolve()) for a in allowed_dirs):
            return {"error": "Seguridad: solo puedo escribir en Downloads, Documents o el proyecto."}
        if mode not in ("w", "a"):
            return {"error": "mode debe ser 'w' (write/overwrite) o 'a' (append)."}
        p.parent.mkdir(parents=True, exist_ok=True)
        try:
            with open(p, mode, encoding="utf-8") as f:
                f.write(content)
            return {"written": str(p), "bytes": len(content.encode("utf-8")), "mode": mode}
        except Exception as e:
            return {"error": str(e)}

    def _tool_code_execute(self, code: str, timeout: int = 15) -> dict:
        """[HITL] Execute a Python snippet in an isolated subprocess. Returns stdout/stderr."""
        if len(code) > 8000:
            return {"error": "Code too long (max 8000 chars)."}
        import tempfile, os as _os
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(
                mode="w", suffix=".py", delete=False, encoding="utf-8"
            ) as tmp:
                tmp.write(code)
                tmp_path = tmp.name
            proc = subprocess.run(
                [sys.executable, tmp_path],
                capture_output=True, text=True, timeout=timeout, shell=False,
            )
            return {
                "stdout": proc.stdout[:3000],
                "stderr": proc.stderr[:1000],
                "returncode": proc.returncode,
            }
        except subprocess.TimeoutExpired:
            return {"error": f"Timeout tras {timeout}s de ejecución."}
        except Exception as e:
            return {"error": str(e)}
        finally:
            if tmp_path:
                try:
                    _os.unlink(tmp_path)
                except Exception:
                    pass

    def _tool_http_request(
        self,
        url: str,
        method: str = "GET",
        headers: dict | None = None,
        body: str = "",
        timeout: int = 10,
    ) -> dict:
        """[HITL] Make an HTTP request (GET/POST/PUT/PATCH).

        SSRF-hardened: blocks loopback, RFC1918 private, link-local (incl. cloud
        metadata 169.254.169.254), multicast and reserved targets — including
        hostnames that resolve to them — unless trusted-lab mode is enabled.

        Redirects are followed MANUALLY with a small hop cap so that EVERY hop is
        re-validated by the SSRF guard: a public URL cannot 30x-bounce into an
        internal/metadata address, and a malformed/relative Location fails closed.
        """
        import urllib.parse

        method = method.upper()
        if method not in ("GET", "POST", "PUT", "PATCH", "DELETE", "HEAD"):
            return {"error": f"Método HTTP inválido: {method}"}

        _MAX_REDIRECTS = 5
        _REDIRECT_CODES = (301, 302, 303, 307, 308)

        current_url = url
        cur_method = method
        cur_body = body
        try:
            for _hop in range(_MAX_REDIRECTS + 1):
                # Re-validate the *current* target BEFORE fetching it, so every
                # redirect hop is SSRF-checked, not just the initial URL.
                block_reason = _http_target_blocked(current_url)
                if block_reason:
                    logger.warning(f"http_request bloqueado: {current_url!r} — {block_reason}")
                    return {"error": block_reason}

                resp = requests.request(
                    cur_method, current_url,
                    headers=headers or {},
                    data=cur_body.encode("utf-8") if cur_body else None,
                    timeout=timeout,
                    allow_redirects=False,
                )

                if resp.status_code not in _REDIRECT_CODES:
                    return {
                        "status_code": resp.status_code,
                        "url": str(resp.url),
                        "headers": dict(resp.headers),
                        "body": resp.text[:4000],
                        "encoding": resp.encoding,
                    }

                location = resp.headers.get("Location") or resp.headers.get("location")
                if not location or not str(location).strip():
                    return {"error": "Redirección sin cabecera Location válida (bloqueado)."}

                # Resolve relative redirects against the current URL, then loop so
                # the new target is SSRF-checked before it is ever fetched.
                current_url = urllib.parse.urljoin(current_url, str(location).strip())
                # Browser/requests semantics: 301/302/303 downgrade to a bodyless
                # GET; 307/308 preserve method and body.
                if resp.status_code in (301, 302, 303) and cur_method not in ("GET", "HEAD"):
                    cur_method = "GET"
                    cur_body = ""

            return {"error": f"Demasiadas redirecciones (>{_MAX_REDIRECTS}) — bloqueado."}
        except Exception as e:
            return {"error": str(e)}

    def _tool_decode_payload(self, payload: str, encoding: str = "auto") -> dict:
        """[EXEMPT] Decode: base64, hex, url, rot13, jwt, or auto-detect all."""
        import base64, urllib.parse, codecs

        if len(payload) > 50_000:
            return {"error": "Payload demasiado grande (máx 50k chars)."}

        target = payload.strip()

        def try_b64(s: str) -> str | None:
            try:
                return base64.b64decode(s + "=" * (-len(s) % 4)).decode("utf-8", errors="replace")
            except Exception:
                return None

        def try_hex(s: str) -> str | None:
            try:
                clean = s.replace(" ", "").replace("0x", "").replace("\\x", "")
                return bytes.fromhex(clean).decode("utf-8", errors="replace")
            except Exception:
                return None

        def try_url(s: str) -> str | None:
            try:
                dec = urllib.parse.unquote(s)
                return dec if dec != s else None
            except Exception:
                return None

        def try_jwt(s: str) -> dict | None:
            parts = s.split(".")
            if len(parts) != 3:
                return None
            try:
                h = try_b64(parts[0]) or ""
                p = try_b64(parts[1]) or ""
                return {"header": h, "payload": p, "sig_prefix": parts[2][:16] + "…"}
            except Exception:
                return None

        results: dict = {}
        if encoding == "auto":
            b64 = try_b64(target)
            if b64:
                results["base64"] = b64
            hx = try_hex(target)
            if hx:
                results["hex"] = hx
            ur = try_url(target)
            if ur:
                results["url"] = ur
            jw = try_jwt(target)
            if jw:
                results["jwt"] = jw
            results["rot13"] = codecs.encode(target, "rot_13")
        elif encoding == "base64":
            results["base64"] = try_b64(target) or "Error decodificando base64"
        elif encoding == "hex":
            results["hex"] = try_hex(target) or "Error decodificando hex"
        elif encoding == "url":
            results["url"] = try_url(target) or target
        elif encoding == "rot13":
            results["rot13"] = codecs.encode(target, "rot_13")
        elif encoding == "jwt":
            r = try_jwt(target)
            results["jwt"] = r if r else "No es un JWT válido (formato: header.payload.sig)"
        else:
            return {"error": f"Encoding desconocido: {encoding}. Usa: auto|base64|hex|url|rot13|jwt"}

        if not results:
            return {"error": "No se pudo decodificar — payload inválido para todos los esquemas."}
        results["original_length"] = len(payload)
        return results

    def _tool_hash_file(self, path: str, algorithms: list | None = None) -> dict:
        """[EXEMPT] Compute MD5/SHA1/SHA256/SHA512 hashes of a file."""
        import hashlib
        p = Path(path).expanduser().resolve()
        if not p.exists():
            return {"error": f"Archivo no encontrado: {path}"}
        algos = [a.lower() for a in (algorithms or ["md5", "sha1", "sha256"])
                 if a.lower() in {"md5", "sha1", "sha256", "sha512"}]
        if not algos:
            return {"error": "Algoritmos inválidos. Opciones: md5, sha1, sha256, sha512"}
        try:
            hashers = {a: hashlib.new(a) for a in algos}
            with open(p, "rb") as f:
                while chunk := f.read(65536):
                    for h in hashers.values():
                        h.update(chunk)
            return {
                "file": p.name,
                "size_bytes": p.stat().st_size,
                "hashes": {a: h.hexdigest() for a, h in hashers.items()},
            }
        except Exception as e:
            return {"error": str(e)}

    def _tool_save_note(self, title: str, content: str, tags: list | None = None) -> dict:
        """[EXEMPT] Persist a markdown note to brain/notes.md."""
        from core.memory_router import redact_secrets
        notes_path = Path(__file__).parent.parent / "brain" / "notes.md"
        notes_path.parent.mkdir(parents=True, exist_ok=True)
        now = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
        tag_line = f"**Tags:** {', '.join(tags)}\n" if tags else ""
        entry = f"\n## {redact_secrets(title)}\n*{now}*\n{tag_line}\n{redact_secrets(content)}\n\n---\n"
        try:
            with open(notes_path, "a", encoding="utf-8") as f:
                f.write(entry)
            return {"saved": True, "path": str(notes_path), "bytes": len(entry)}
        except Exception as e:
            return {"error": str(e)}

    def _tool_list_notes(self, query: str = "", limit: int = 10) -> dict:
        """[EXEMPT] List recent notes from brain/notes.md, optionally filtered."""
        notes_path = Path(__file__).parent.parent / "brain" / "notes.md"
        if not notes_path.exists():
            return {"notes": [], "message": "No hay notas guardadas aún."}
        try:
            raw = notes_path.read_text(encoding="utf-8")
            entries = [e.strip() for e in raw.split("---") if e.strip()]
            if query:
                q = query.lower()
                entries = [e for e in entries if q in e.lower()]
            entries = entries[-limit:]
            return {"notes": entries, "total_shown": len(entries)}
        except Exception as e:
            return {"error": str(e)}

    def _tool_project_note(self, kind: str, text: str) -> dict:
        """[EXEMPT] Record a project fact (goal/decision/task/blocked/question/artifact).

        V63 M8 — persisted via the memory fabric at scope=project with provenance
        + timestamp, so JARVIS stays aware of ongoing work. Writes only to
        JARVIS's own memory (LOW_IMPACT, non-HITL — same tier as save_note).
        """
        import asyncio as _asyncio
        from core.project_context import ProjectFactType, record_project_fact
        valid = {t.value for t in ProjectFactType}
        k = (kind or "").strip().lower()
        if k not in valid:
            return {"error": f"kind must be one of {sorted(valid)}"}
        if not (text or "").strip():
            return {"error": "text is required"}
        try:
            ok = _asyncio.run(record_project_fact(ProjectFactType(k), text))
            return {"recorded": bool(ok), "kind": k}
        except Exception as e:
            return {"error": str(e)}

    def _tool_project_status(self, query: str = "") -> dict:
        """[EXEMPT] Recall current project context — goals, decisions, tasks,
        blockers, open questions, artifacts — grouped by type (READ_ONLY)."""
        import asyncio as _asyncio
        from core.project_context import summarize_project
        try:
            return _asyncio.run(summarize_project(query))
        except Exception as e:
            return {"error": str(e)}

    def _tool_git_query(self, operation: str = "status", args: str = "") -> dict:
        """[EXEMPT] Read-only git: status, diff, log, show, branch, stash."""
        allowed_ops = {"status", "diff", "log", "show", "branch", "stash"}
        if operation not in allowed_ops:
            return {"error": f"Operación inválida. Permitidas: {', '.join(sorted(allowed_ops))}"}
        if _FORBIDDEN_CHARS_RE.search(args):
            return {"error": "args contiene metacaracteres prohibidos."}

        argv = ["git", operation]
        if args.strip():
            try:
                extra = shlex.split(args)
            except ValueError as e:
                return {"error": f"args malformados: {e}"}
            write_flags = {"--add", "-A", "--amend", "-m", "--force", "-f", "--delete", "-d"}
            if any(f in write_flags for f in extra):
                return {"error": "Flags de escritura no permitidos en modo read-only."}
            argv.extend(extra)

        if operation == "log" and "--oneline" not in argv:
            argv.append("--oneline")
        if operation == "log" and not any(a.startswith("-n") for a in argv):
            argv += ["-n", "20"]

        try:
            proc = subprocess.run(argv, capture_output=True, text=True, timeout=15, shell=False)
            return {
                "operation": operation,
                "stdout": proc.stdout[:3000],
                "stderr": proc.stderr[:500],
                "returncode": proc.returncode,
            }
        except FileNotFoundError:
            return {"error": "git no encontrado en el PATH del sistema."}
        except subprocess.TimeoutExpired:
            return {"error": "git timeout (15s)."}
        except Exception as e:
            return {"error": str(e)}

    def _tool_port_lookup(self, port: int | str, protocol: str = "tcp") -> dict:
        """[EXEMPT] Resolve port number → service name + risk level."""
        import socket as _sock
        _WELL_KNOWN: dict[int, str] = {
            20: "FTP-Data", 21: "FTP", 22: "SSH", 23: "Telnet", 25: "SMTP",
            53: "DNS", 67: "DHCP-Server", 68: "DHCP-Client", 69: "TFTP",
            80: "HTTP", 110: "POP3", 111: "RPCBind", 123: "NTP",
            135: "MS-RPC", 137: "NetBIOS-NS", 139: "NetBIOS-SS",
            143: "IMAP", 161: "SNMP", 162: "SNMP-Trap",
            389: "LDAP", 443: "HTTPS", 445: "SMB", 465: "SMTPS",
            500: "IKE/IPsec", 514: "Syslog", 587: "SMTP-Submission",
            636: "LDAPS", 993: "IMAPS", 995: "POP3S",
            1080: "SOCKS", 1194: "OpenVPN", 1433: "MSSQL", 1521: "Oracle",
            2049: "NFS", 2375: "Docker", 2376: "Docker-TLS",
            3306: "MySQL", 3389: "RDP", 4444: "Metasploit/C2",
            5432: "PostgreSQL", 5900: "VNC", 5985: "WinRM-HTTP",
            5986: "WinRM-HTTPS", 6379: "Redis", 6443: "Kubernetes-API",
            8080: "HTTP-Alt", 8443: "HTTPS-Alt", 8888: "Jupyter/Proxy",
            9200: "Elasticsearch", 11211: "Memcached", 27017: "MongoDB",
        }
        _HIGH_RISK = {4444, 31337, 1337, 5900, 3389, 445, 23, 135, 139, 161}
        _MED_RISK  = {21, 22, 25, 80, 110, 143, 389, 636, 2375, 6379, 9200, 11211, 27017}

        try:
            p = int(port)
        except (ValueError, TypeError):
            return {"error": "port debe ser un número entero."}
        if not 0 <= p <= 65535:
            return {"error": "port fuera de rango (0-65535)."}

        name = _WELL_KNOWN.get(p)
        if not name:
            try:
                name = _sock.getservbyport(p, protocol.lower())
            except OSError:
                name = "Unknown"

        risk = "HIGH" if p in _HIGH_RISK else ("MEDIUM" if p in _MED_RISK else "LOW")
        return {"port": p, "protocol": protocol.lower(), "service": name, "risk_level": risk}

    def _tool_regex_test(self, pattern: str, text: str, flags: str = "") -> dict:
        """[EXEMPT] Test a regex against text. Returns all matches with groups."""
        import re as _re
        flag_map = {"i": _re.IGNORECASE, "m": _re.MULTILINE, "s": _re.DOTALL}
        f_val = 0
        for ch in flags.lower():
            if ch in flag_map:
                f_val |= flag_map[ch]
        try:
            rx = _re.compile(pattern, f_val)
        except _re.error as e:
            return {"error": f"Regex inválido: {e}"}
        matches = list(rx.finditer(text))
        return {
            "pattern": pattern,
            "flags": flags or "none",
            "total_matches": len(matches),
            "is_match": bool(matches),
            "matches": [
                {"match": m.group(0), "start": m.start(), "end": m.end(), "groups": list(m.groups())}
                for m in matches[:25]
            ],
        }


# ── RedTeamShellExecutor ──────────────────────────────────────────────────────

class RedTeamShellExecutor:
    """Permissive shell executor for Red Team operator use.

    Sits alongside ToolExecutor without modifying it.
    Authorization gate: dynamic trust model + NATO OTP via ToolExecutor._challenge().

    Security layers:
      Layer 0 — YARA scan of the raw command string.
      Hard filter — _CRITICAL_BLOCK patterns are unconditional; no OTP override.
      Layer 1 — _classify(): binary must exist on OS; unlisted binaries escalate.
      Layer 2 — Dynamic trust challenge (NONE / CONFIRM / FULL_NATO).
      Layer 3 — subprocess.run with shell=False always.
    Full audit trail via ToolExecutor._audit.log_action().
    """

    def __init__(self, tool_executor: "ToolExecutor") -> None:
        self._te = tool_executor  # borrow _challenge() and _audit from ToolExecutor
        self._trust_profile: dict = {}
        self._session_commands: list[str] = []
        self._profile_loaded: bool = False

    def _classify(self, command: str, yara_hits: list) -> tuple[str, bool]:
        tokens = shlex.split(command, posix=False)
        if not tokens:
            raise ValueError("Empty command")

        binary = Path(tokens[0]).name.lower().removesuffix(".exe")

        if shutil.which(binary) is None:
            raise ValueError(f"[BLOCKED] Binary '{binary}' not found on OS")

        in_allowlist = binary in COMMAND_ALLOWLIST
        suspicious   = bool(yara_hits) or any(
            p in command.lower() for p in _SUSPICIOUS_PATTERNS
        )

        if not in_allowlist:
            return ("unlisted_escalation", True)
        return ("allowlisted", suspicious)

    async def execute_shell(self, command: str, reasoning: str = "") -> dict:
        result: dict = {
            "command":    command,
            "authorized": False,
            "stdout":     "",
            "stderr":     "",
            "error":      None,
        }
        binary_status = "unknown"
        auth_word     = ""

        # Lazy-load trust profile on first execution
        if not self._profile_loaded:
            from core.trust_engine import load_profile
            self._trust_profile = await load_profile()
            self._profile_loaded = True

        try:
            # Layer 0 — YARA scan (scan_command expects list[str])
            from core.yara_analyzer import scan_command
            yara_hits = await scan_command(shlex.split(command))

            # Hard filter — unconditional, no OTP override
            cmd_low = command.lower()
            for pattern in _CRITICAL_BLOCK:
                if pattern in cmd_low:
                    raise ValueError(f"[HARD BLOCK] OS-destructive pattern: '{pattern}'")

            # F3 — strict allowlist validation (same pipeline as run_shell_command).
            # Off-allowlist binaries, shell metacharacters, encoded/inline-exec
            # flags and system-path arguments are REFUSED here, BEFORE any trust
            # scoring — so a high trust score (or a denylist gap like the
            # double-space "rm  -rf  /") can never authorize an unlisted binary.
            # Lab-only binaries pass solely via the explicit lab allowlist under
            # trusted-lab mode; _classify then marks them unlisted → FULL_NATO.
            extra = _LAB_COMMAND_ALLOWLIST if _trusted_lab_enabled() else frozenset()
            ok, verr, _ = _validate_command(command, extra)
            if not ok:
                raise ValueError(f"[BLOCKED] {verr}")

            # Layer 1 — classify binary
            binary_status, _ = self._classify(command, yara_hits)
            tokens = shlex.split(command, posix=False)
            binary = Path(tokens[0]).name.lower().removesuffix(".exe")

            # Layer 2 — dynamic trust challenge
            from core.trust_engine import (
                get_challenge_level, ChallengeLevel, update_profile,
            )
            level, score = get_challenge_level(
                binary, command, binary_status, yara_hits,
                self._trust_profile, self._session_commands,
            )

            # Broadcast trust decision to AURA
            await _aura_broadcast({
                "type":      "trust_decision",
                "binary":    binary,
                "level":     level.value,
                "score":     round(score, 2),
                "timestamp": datetime.now(timezone.utc).isoformat(),
            })

            loop = asyncio.get_running_loop()
            if level == ChallengeLevel.CONFIRM:
                response = await loop.run_in_executor(
                    None,
                    lambda: input("[CONFIRM] Type 'yes' to proceed: ").strip().lower(),
                )
                if response not in ("yes", "y"):
                    raise ValueError("[DENIED] Operator declined confirmation")
            elif level == ChallengeLevel.FULL_NATO:
                auth_ok, auth_word = await self._te._challenge(
                    tool_name="run_shell_command",
                    preview=command[:120],
                )
                if not auth_ok:
                    raise ValueError(f"[DENIED] NATO challenge failed: {auth_word}")

            # Layer 3 — execute, shell=False always
            result["authorized"] = True
            proc = await loop.run_in_executor(
                None,
                lambda: subprocess.run(
                    shlex.split(command),
                    shell=False,
                    capture_output=True,
                    text=True,
                    timeout=60,
                ),
            )
            result["stdout"] = proc.stdout
            result["stderr"] = proc.stderr

            # Update trust profile on successful execution
            asyncio.create_task(update_profile(binary, self._trust_profile))
            self._session_commands.append(binary)

        except ValueError as e:
            result["error"] = str(e)

        finally:
            self._te._audit.log_action(
                "run_shell_command",
                reasoning,
                auth_word,
                "authorized" if result["authorized"] else "denied",
                f"RedTeamShell | binary_status={binary_status} | cmd={command[:80]}",
            )

        return result
